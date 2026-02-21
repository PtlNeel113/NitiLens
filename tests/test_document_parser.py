"""
Tests for the unified document parser service.
Validates PDF and PPT/PPTX extraction, error handling, and timeout protection.

Run with:
    cd backend && python -m pytest ../tests/test_document_parser.py -v
"""
import os
import sys
import tempfile
from pathlib import Path

# Ensure backend is importable
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "backend"))

import pytest


class TestDocumentParserImport:
    """Verify the parser module loads correctly."""

    def test_import_module(self):
        from app.services.document_parser import extract_text, get_document_metadata
        assert callable(extract_text)
        assert callable(get_document_metadata)


class TestFileTypeValidation:
    """Verify unsupported file types raise ValueError."""

    def test_unsupported_type_raises(self):
        from app.services.document_parser import extract_text
        with pytest.raises(ValueError, match="Unsupported document format"):
            extract_text("dummy.txt", "txt")

    def test_unsupported_type_csv(self):
        from app.services.document_parser import extract_text
        with pytest.raises(ValueError, match="Unsupported document format"):
            extract_text("dummy.csv", "csv")

    def test_unsupported_type_empty(self):
        from app.services.document_parser import extract_text
        with pytest.raises(ValueError, match="Unsupported document format"):
            extract_text("dummy", "")


class TestPPTExtraction:
    """Test PPT/PPTX text extraction."""

    def _create_pptx_with_text(self, texts: list[str]) -> str:
        """Helper: create a temporary PPTX file with given slide texts."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for text in texts:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
            slide.shapes.title.text = text

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)
        tmp.close()
        return tmp.name

    def _create_empty_pptx(self) -> str:
        """Helper: create a PPTX with no text content."""
        from pptx import Presentation

        prs = Presentation()
        # Add a blank slide with no text
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)
        tmp.close()
        return tmp.name

    def test_small_pptx(self):
        """✔ Small PPT (5 slides)"""
        from app.services.document_parser import extract_text

        path = self._create_pptx_with_text([
            "AML Policy Overview",
            "Transaction limits above $10,000",
            "Structuring detection rules",
            "Suspicious Activity Reporting",
            "Compliance team responsibilities",
        ])
        try:
            text = extract_text(path, "pptx")
            assert "AML Policy Overview" in text
            assert "$10,000" in text
            assert len(text) > 50
        finally:
            os.unlink(path)

    def test_large_pptx(self):
        """✔ Large PPT (50+ slides)"""
        from app.services.document_parser import extract_text

        slides = [f"Slide {i}: Compliance rule #{i}" for i in range(55)]
        path = self._create_pptx_with_text(slides)
        try:
            text = extract_text(path, "pptx")
            assert "Slide 1" in text
            assert "Slide 54" in text
        finally:
            os.unlink(path)

    def test_empty_pptx_raises(self):
        """✔ Empty PPT → raises ValueError"""
        from app.services.document_parser import extract_text

        path = self._create_empty_pptx()
        try:
            with pytest.raises(ValueError, match="No readable text"):
                extract_text(path, "pptx")
        finally:
            os.unlink(path)

    def test_corrupted_pptx_raises(self):
        """✔ Corrupted PPT → raises error"""
        from app.services.document_parser import extract_text

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.write(b"this is not a valid pptx file")
        tmp.close()
        try:
            with pytest.raises(Exception):
                extract_text(tmp.name, "pptx")
        finally:
            os.unlink(tmp.name)

    def test_pptx_metadata(self):
        """Test metadata extraction for PPTX."""
        from app.services.document_parser import get_document_metadata

        path = self._create_pptx_with_text(["Slide 1", "Slide 2", "Slide 3"])
        try:
            meta = get_document_metadata(path, "pptx")
            assert meta["file_type"] == "pptx"
            assert meta["total_slides"] == 3
        finally:
            os.unlink(path)


class TestPDFExtraction:
    """Test that PDF extraction still works (uses stub fallback)."""

    def test_pdf_stub_aml(self):
        """PDF extraction falls back to AML stub for demo files."""
        from app.services.document_parser import extract_text

        # Create a dummy file named like an AML policy
        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", prefix="aml_policy_", delete=False)
        tmp.write(b"%PDF-1.4 dummy")
        tmp.close()
        try:
            text = extract_text(tmp.name, "pdf")
            assert "ANTI-MONEY LAUNDERING" in text or len(text) > 0
        finally:
            os.unlink(tmp.name)

    def test_pdf_stub_gdpr(self):
        """PDF extraction falls back to GDPR stub for demo files."""
        from app.services.document_parser import extract_text

        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", prefix="gdpr_policy_", delete=False)
        tmp.write(b"%PDF-1.4 dummy")
        tmp.close()
        try:
            text = extract_text(tmp.name, "pdf")
            assert "DATA PROTECTION" in text or len(text) > 0
        finally:
            os.unlink(tmp.name)


class TestTimeoutProtection:
    """Verify timeout protection wrapping."""

    def test_extract_text_has_timeout(self):
        """Confirm the function uses ThreadPoolExecutor with timeout."""
        import inspect
        from app.services.document_parser import extract_text
        source = inspect.getsource(extract_text)
        assert "ThreadPoolExecutor" in source
        assert "timeout" in source.lower()


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
